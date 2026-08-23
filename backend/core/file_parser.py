# -*- coding: utf-8 -*-
"""文件解析（照 DeepTutor 可插拔解析引擎）：优先 markitdown 统一转 Markdown，
失败 fallback 到基础解析（pypdf/python-docx/python-pptx/纯文本）。

markitdown（微软开源，pip install 'markitdown[pdf,docx,pptx,xlsx]'）
支持：pdf/docx/pptx/xlsx/xls/html/csv/json/xml/txt/md/epub/png/jpg/jpeg/gif/webp。
"""

TEXT_EXTS = {"txt", "md", "py", "js", "ts", "json", "csv", "html", "css", "log", "yaml", "yml"}

# markitdown 支持的扩展名（DeepTutor markitdown engine _SUPPORTED）
MARKITDOWN_EXTS = {"pdf", "docx", "pptx", "xlsx", "xls", "html", "htm", "csv", "json", "xml", "txt", "md", "epub", "png", "jpg", "jpeg", "gif", "webp"}


def _parse_with_markitdown(data: bytes) -> str | None:
    """markitdown 转 Markdown；未安装或失败返回 None（调用方 fallback）"""
    try:
        from markitdown import MarkItDown
        import io as _io
        md = MarkItDown()
        result = md.convert(_io.BytesIO(data))
        text = (result.text_content or "").strip()
        return text or None
    except Exception:
        return None


def _parse_fallback(filename: str, data: bytes) -> str:
    """基础解析（markitdown 不可用时）"""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        import io as _io
        from pypdf import PdfReader
        reader = PdfReader(_io.BytesIO(data))
        parts = []
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            if t:
                parts.append(t)
        return "\n\n".join(parts)
    if ext == "docx":
        import io as _io
        from docx import Document
        doc = Document(_io.BytesIO(data))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    if ext == "pptx":
        import io as _io
        from pptx import Presentation
        prs = Presentation(_io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_parts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if slide_parts:
                parts.append("第" + str(i + 1) + "页: " + " / ".join(slide_parts))
        return "\n".join(parts)
    if ext in TEXT_EXTS:
        return data.decode("utf-8", errors="replace")
    return data.decode("utf-8", errors="replace")


def parse_file(filename: str, data: bytes) -> str:
    """解析文件内容为文本：markitdown 优先（pdf/docx/pptx/xlsx/html 等），失败 fallback 基础解析。"""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in MARKITDOWN_EXTS:
        text = _parse_with_markitdown(data)
        if text:
            return text
    # fallback
    try:
        return _parse_fallback(filename, data)
    except Exception:
        return ""
